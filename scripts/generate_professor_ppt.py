import json
import os
from datetime import datetime

from pptx import Presentation
from pptx.util import Inches, Pt


BASE_DIR = os.path.abspath(os.path.join(os.path.dirname(__file__), ".."))
METRICS_PATH = os.path.join(BASE_DIR, "reports", "unsw_evaluation", "metrics_test.json")
MODEL_INFO_PATH = os.path.join(BASE_DIR, "models", "unsw_metrics.json")
REPORTS_DIR = os.path.join(BASE_DIR, "reports", "unsw_evaluation")
OUTPUT_PPT = os.path.join(BASE_DIR, "Professor_Review_IDS.pptx")


def add_title_slide(prs, title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle


def add_bullet_slide(prs, title, bullets):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    tf = slide.placeholders[1].text_frame
    tf.clear()

    for i, text in enumerate(bullets):
        if i == 0:
            tf.text = text
            tf.paragraphs[0].font.size = Pt(20)
        else:
            p = tf.add_paragraph()
            p.text = text
            p.font.size = Pt(20)


def add_two_column_slide(prs, title, left_title, left_lines, right_title, right_lines):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = title

    left = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(6.1), Inches(5.8)).text_frame
    right = slide.shapes.add_textbox(Inches(6.9), Inches(1.2), Inches(6.1), Inches(5.8)).text_frame

    left.text = left_title
    left.paragraphs[0].font.bold = True
    left.paragraphs[0].font.size = Pt(18)

    for line in left_lines:
        p = left.add_paragraph()
        p.text = line
        p.font.size = Pt(14)

    right.text = right_title
    right.paragraphs[0].font.bold = True
    right.paragraphs[0].font.size = Pt(18)

    for line in right_lines:
        p = right.add_paragraph()
        p.text = line
        p.font.size = Pt(14)


def add_metrics_table_slide(prs, title, metrics):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = title

    rows, cols = 8, 6
    table_shape = slide.shapes.add_table(rows, cols, Inches(0.5), Inches(1.3), Inches(12.3), Inches(4.5))
    table = table_shape.table

    headers = ["Model", "Accuracy", "Precision", "Recall", "F1", "ROC-AUC"]
    for c, h in enumerate(headers):
        table.cell(0, c).text = h
        for p in table.cell(0, c).text_frame.paragraphs:
            p.font.bold = True
            p.font.size = Pt(12)

    order = ["rf", "xgb", "lgb", "et", "mlp", "ensemble_raw", "ensemble"]
    names = {
        "rf": "Random Forest",
        "xgb": "XGBoost",
        "lgb": "LightGBM",
        "et": "Extra Trees",
        "mlp": "MLP",
        "ensemble_raw": "Ensemble Raw",
        "ensemble": "Ensemble Calibrated",
    }

    for r, key in enumerate(order, start=1):
        m = metrics.get(key, {})
        row = [
            names[key],
            f"{m.get('accuracy', 0):.4f}",
            f"{m.get('precision', 0):.4f}",
            f"{m.get('recall', 0):.4f}",
            f"{m.get('f1', 0):.4f}",
            f"{m.get('roc_auc', 0):.4f}",
        ]
        for c, value in enumerate(row):
            table.cell(r, c).text = value
            for p in table.cell(r, c).text_frame.paragraphs:
                p.font.size = Pt(11)


def add_image_slide(prs, title, image_file, caption):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = title

    if os.path.exists(image_file):
        slide.shapes.add_picture(image_file, Inches(0.7), Inches(1.2), width=Inches(11.8))

    tb = slide.shapes.add_textbox(Inches(0.7), Inches(6.8), Inches(11.8), Inches(0.5))
    tf = tb.text_frame
    tf.text = caption
    tf.paragraphs[0].font.size = Pt(13)


def main():
    with open(METRICS_PATH, "r", encoding="utf-8") as f:
        metrics = json.load(f)

    with open(MODEL_INFO_PATH, "r", encoding="utf-8") as f:
        model_info = json.load(f)

    prs = Presentation()

    add_title_slide(
        prs,
        "Network Intrusion Detection System",
        f"Professor Review | UNSW-NB15 | {datetime.now().strftime('%Y-%m-%d %H:%M')}"
    )

    add_bullet_slide(prs, "Project Intro", [
        "Goal: classify network traffic as Normal or Attack",
        "Dataset: UNSW-NB15 train and test split",
        "System: Flask backend + React dashboard",
        "Approach: stacked ensemble + probability calibration",
        "Output: prediction, attack probability, threat level"
    ])

    add_two_column_slide(
        prs,
        "Key Parameters The Model Checks",
        "Traffic and Packet Behavior",
        [
            "dur: connection duration",
            "spkts, dpkts: packet counts",
            "sbytes, dbytes: byte counts",
            "rate: traffic rate",
            "sinpkt, dinpkt: packet timing",
            "sjit, djit: jitter",
            "sttl, dttl: TTL values",
            "sload, dload: traffic load"
        ],
        "Protocol and Pattern Indicators",
        [
            "proto, service, state: protocol/service/state",
            "tcprtt, synack, ackdat: TCP timing",
            "ct_srv_src: repeated service from source",
            "ct_dst_ltm: repeated destination access",
            "ct_src_dport_ltm: repeated dst port from source",
            "ct_flw_http_mthd: HTTP method count",
            "is_ftp_login: FTP login indicator",
            "is_sm_ips_ports: src/dst same IP-port pattern"
        ]
    )

    add_bullet_slide(prs, "Models and Ensemble", [
        "Base models: Random Forest, XGBoost, LightGBM, Extra Trees, MLP",
        "Each base model outputs attack probability",
        "Meta model combines base probabilities (stacking)",
        f"Meta model used: {model_info.get('stacking', {}).get('meta_model', 'LogisticRegression')}",
        "Final probability is calibrated for interpretability"
    ])

    add_metrics_table_slide(prs, "Test Metrics (Individual + Ensemble)", metrics)

    add_image_slide(
        prs,
        "Confusion Matrix Overview",
        os.path.join(REPORTS_DIR, "confusion_matrices_test.png"),
        "Model-wise confusion matrices on UNSW test set"
    )

    add_image_slide(
        prs,
        "ROC Curve Comparison",
        os.path.join(REPORTS_DIR, "roc_curves_test.png"),
        "AUC comparison across models"
    )

    add_image_slide(
        prs,
        "Calibration Curve",
        os.path.join(REPORTS_DIR, "ensemble_calibration_curve_test.png"),
        "Predicted probabilities vs observed outcomes"
    )

    add_bullet_slide(prs, "Conclusion", [
        "All individual models are working and contribute to final output",
        "Ensemble improves stability and overall performance",
        f"Ensemble test accuracy: {metrics.get('ensemble', {}).get('accuracy', 0):.4f}",
        "Model is ready for demonstration with live dashboard"
    ])

    prs.save(OUTPUT_PPT)
    print(f"Created: {OUTPUT_PPT}")


if __name__ == "__main__":
    main()
