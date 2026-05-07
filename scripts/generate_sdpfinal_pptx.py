from pathlib import Path
import json

try:
    from pptx import Presentation
except Exception:
    import subprocess
    import sys
    subprocess.check_call([sys.executable, '-m', 'pip', 'install', 'python-pptx'])
    from pptx import Presentation

WORKSPACE = Path(__file__).resolve().parents[1]
TEMPLATE = WORKSPACE / 'SCOPE_SDP_Template_May_2026.pptx'
OUTPUT = WORKSPACE / 'sdpfinalreview.pptx'

README = WORKSPACE / 'README.md'
MODEL_JSON = WORKSPACE / 'models' / 'unsw_metrics.json'
TRAIN_JSON = WORKSPACE / 'reports' / 'unsw_evaluation' / 'metrics_train.json'
TEST_JSON = WORKSPACE / 'reports' / 'unsw_evaluation' / 'metrics_test.json'
QUANT_JSON = WORKSPACE / 'reports' / 'unsw_evaluation' / 'ensemble_probability_quantiles.json'
CALIB_JSON = WORKSPACE / 'reports' / 'unsw_evaluation' / 'calibration_retry_comparison.json'


def load_json(path):
    if path.exists():
        with path.open('r', encoding='utf-8') as f:
            return json.load(f)
    return {}


def read_text(path):
    if path.exists():
        return path.read_text(encoding='utf-8')
    return ''


def set_text(shape, lines):
    text_frame = shape.text_frame
    text_frame.clear()
    if not lines:
        return
    first = text_frame.paragraphs[0]
    first.text = str(lines[0])
    for line in lines[1:]:
        paragraph = text_frame.add_paragraph()
        paragraph.text = str(line)


def set_bullets(shape, bullets):
    text_frame = shape.text_frame
    text_frame.clear()
    for index, bullet in enumerate(bullets):
        paragraph = text_frame.paragraphs[0] if index == 0 else text_frame.add_paragraph()
        paragraph.text = str(bullet)
        paragraph.level = 0


def shape_text(shape):
    if not hasattr(shape, 'text'):
        return ''
    return shape.text.strip().lower()


readme = read_text(README)
model = load_json(MODEL_JSON)
train = load_json(TRAIN_JSON)
test = load_json(TEST_JSON)
quant = load_json(QUANT_JSON)
calib = load_json(CALIB_JSON)

prs = Presentation(str(TEMPLATE))

# Slide 1: title page
slide = prs.slides[0]
set_text(slide.shapes[0], [
    'Senior Design Project Review',
    'Network Intrusion Detection System'
])
set_text(slide.shapes[1], ['SDP ID:', 'sdpfinalreview'])
set_text(slide.shapes[2], ['Presented By:', 'sdpfinalreview'])
set_text(slide.shapes[3], ['SCOPE', 'VIT-AP University,', 'Amravati, India'])
# keep slide number and decorative shapes intact

# Slide 2: outline
set_bullets(prs.slides[1].shapes[1], [
    'Introduction',
    'Project Overview and Problem Statement',
    'Motivations',
    'Background & Related Work / Literature Review',
    'Project Objectives',
    'Proposed Solution',
    'Simulation and Results',
    'Summary'
])

# Slide 3: introduction
set_bullets(prs.slides[2].shapes[1], [
    'Project Overview',
    'AI-powered network intrusion detection system for UNSW-NB15 traffic.',
    'Uses an ensemble of machine-learning models to classify traffic as normal or attack.',
    'Exposes a Flask REST API and a React dashboard for real-time and batch analysis.',
    'Problem Statement',
    'Manual monitoring is too slow for modern network traffic volumes.',
    'A practical IDS must be accurate, fast, explainable, and usable by operators.'
])

# Slide 4: motivations
set_bullets(prs.slides[3].shapes[1], [
    'Detect attacks earlier and reduce response time for suspicious traffic.',
    'Provide probabilistic threat scores instead of a hard yes/no output only.',
    'Support both single-sample and batch prediction workflows.',
    'Show results in a modern UI so the system can be demonstrated and reviewed easily.',
    'Use the UNSW-NB15 dataset because it is widely used for network intrusion research.'
])

# Slide 5: literature review
set_bullets(prs.slides[4].shapes[1], [
    'Signature-based IDS tools are effective for known attacks but weaker on new variants.',
    'Traditional ML-based IDS pipelines often improve detection by learning from network features.',
    'Random Forest, gradient boosting, and ensemble methods are common baselines for tabular IDS data.',
    'Probability calibration is important because raw model scores can be overconfident.',
    'A web dashboard improves practical usability compared with offline-only evaluation scripts.'
])

# Slide 6: gaps and improvements
set_bullets(prs.slides[5].shapes[1], [
    'Many prior systems stop at model training and do not expose a live prediction interface.',
    'Some models report accuracy but not calibrated probabilities or calibration metrics.',
    'Single-model approaches are less robust than a stacked ensemble on tabular traffic data.',
    'This project adds backend API endpoints, frontend visualization, batch processing, and model metadata.',
    'The final review package also includes reproducible reports and evaluation artifacts.'
])

# Slide 7: objectives
set_bullets(prs.slides[6].shapes[1], [
    'Build a working intrusion detection pipeline for the UNSW-NB15 dataset.',
    'Train and evaluate multiple models and combine them into an ensemble.',
    'Expose prediction, batch prediction, health, and model-info endpoints.',
    'Provide an interactive frontend for real-time threat monitoring.',
    'Document performance, calibration, and implementation outcomes clearly.'
])

# Slide 8: solution architecture
set_bullets(prs.slides[7].shapes[1], [
    'Input data flows from raw feature vectors into preprocessing and scaling.',
    'Base learners produce attack probabilities that are stacked by a logistic meta-model.',
    'A calibration layer adjusts the ensemble probability for better reliability.',
    'The Flask backend serves predictions and model metadata over REST.',
    'The React frontend presents live status, charts, batch summaries, and activity history.'
])

# Slide 9: components
set_bullets(prs.slides[8].shapes[1], [
    'Backend API: prediction, batch prediction, health, sample batch, and model-info endpoints.',
    'Model artifacts: base estimators, scaler, feature names, categorical mappings, and calibrator.',
    'Frontend dashboard: single prediction form, batch workflow, charts, logs, and summary panels.',
    'Training and evaluation scripts: model training, metrics generation, and report creation.',
    'Artifacts folder: saved metrics, calibration comparison, and probability distribution summaries.'
])

# Slide 10: algorithms and tools
set_bullets(prs.slides[9].shapes[1], [
    'Algorithms: Random Forest, XGBoost, LightGBM, Extra Trees, MLP, and Logistic Regression meta-model.',
    'Calibration: temperature scaling selected for the final UNSW profile.',
    'Data tooling: pandas, numpy, joblib, scikit-learn.',
    'Backend: Flask with Flask-CORS.',
    'Frontend: React, Vite, and Recharts.'
])

# Slide 11: dataset description
set_bullets(prs.slides[10].shapes[1], [
    f'Active dataset/profile: UNSW-NB15',
    f'Training set shape: {model.get("train_shape", [82332, 45])[0]} samples x {model.get("train_shape", [82332, 45])[1]} columns',
    f'Test set shape: {model.get("test_shape", [175341, 45])[0]} samples x {model.get("test_shape", [175341, 45])[1]} columns',
    f'Feature count used by the active model: {model.get("feature_count", 42)}',
    f'Categorical columns handled explicitly: {", ".join(model.get("categorical_columns", []))}',
    'Model output: class label plus calibrated attack probability.'
])

# Slide 12: parameters and initial conditions
stacking = model.get('stacking', {})
calibration = model.get('calibration', {})
set_bullets(prs.slides[11].shapes[1], [
    f'Base model order: {", ".join(stacking.get("base_model_order", []))}',
    f'Meta-model: {stacking.get("meta_model", "LogisticRegression")}',
    f'OOF folds: {stacking.get("oof_folds", 5)}',
    f'Calibration split: {stacking.get("calibration_split", 0.2)}',
    f'Selected calibration method: {calibration.get("selected_method", "temperature")}',
    f'Temperature value: {calibration.get("temperature", 2.5)}',
    'Threat thresholds in the UI: Low < 40%, Medium 40-70%, High > 70%'
])

# Slide 13: outcomes
ensemble = model.get('ensemble_calibrated') or model.get('ensemble_raw') or model.get('ensemble') or {}
metrics_test_ensemble = test.get('ensemble') or test.get('ensemble_raw') or {}
metrics_train_ensemble = train.get('ensemble') or train.get('ensemble_raw') or {}
set_bullets(prs.slides[12].shapes[1], [
    f'Ensemble test accuracy: {metrics_test_ensemble.get("accuracy", ensemble.get("accuracy", "N/A"))}',
    f'Ensemble test precision: {metrics_test_ensemble.get("precision", "N/A")}',
    f'Ensemble test recall: {metrics_test_ensemble.get("recall", "N/A")}',
    f'Ensemble test F1: {metrics_test_ensemble.get("f1", "N/A")}',
    f'Ensemble calibrated Brier score: {ensemble.get("brier", ensemble.get("brier_score", "N/A"))}',
    'Frontend outcomes: live status, trend chart, pie chart, batch summaries, and system logs.'
])

# Slide 14: summary
set_bullets(prs.slides[13].shapes[1], [
    'Key finding: the stacked ensemble gives strong overall detection performance on UNSW-NB15.',
    'Calibration reduces overconfident predictions and improves probability usefulness.',
    'Limitation: the system is model-based and still depends on feature quality and training data coverage.',
    'Future work: live packet capture, richer explainability, retraining automation, and deployment hardening.'
])

# Slide 15: references
set_text(prs.slides[14].shapes[1], [
    'References',
    'Moustafa, N., & Slay, J. (2015). UNSW-NB15: a comprehensive data set for network intrusion detection systems.',
    'Pedregosa, F., et al. (2011). Scikit-learn: Machine Learning in Python. Journal of Machine Learning Research.',
    'Flask Documentation. https://flask.palletsprojects.com/',
    'React Documentation. https://react.dev/',
    'Vite Documentation. https://vitejs.dev/',
    'Recharts Documentation. https://recharts.org/'
])

# Slide 16: closing slide
set_text(prs.slides[15].shapes[0], ['Thank You'])
set_text(prs.slides[15].shapes[1], ['sdpfinalreview'])

prs.save(str(OUTPUT))
print(f'Saved {OUTPUT}')
